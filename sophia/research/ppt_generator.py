try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
"""Academic presentation slide generator (Phase K).

K-1: PPT structure planner
K-2: HTML slide renderer (no python-pptx per user requirement)
"""

from __future__ import annotations

import html
from typing import Any, Dict, List, Optional


class SlidePlanner:
    """Plans academic presentation slide structures."""

    # Predefined slide templates for conference mode
    _CONFERENCE_TEMPLATE = [
        {"type": "title", "title": "标题页", "source": "论文标题", "notes": "包含论文标题、作者、单位、会议名称和日期"},
        {"type": "content", "title": "研究背景与问题", "source": "引言", "notes": "用1-2句话说明研究背景，清晰提出研究问题或核心议题"},
        {"type": "content", "title": "文献回顾", "source": "文献综述", "notes": "概括2-3个关键研究脉络，指出研究空白"},
        {"type": "content", "title": "理论框架", "source": "理论框架", "notes": "展示核心理论模型或分析框架，可用图示"},
        {"type": "content", "title": "研究问题/假设", "source": "研究设计", "notes": "明确列出研究问题或研究假设"},
        {"type": "content", "title": "研究方法", "source": "研究方法", "notes": "简要说明研究设计、数据来源、分析方法"},
        {"type": "content", "title": "数据来源与样本", "source": "研究方法/数据", "notes": "展示样本特征、数据收集方式、时间范围"},
        {"type": "content", "title": "主要发现（一）", "source": "研究结果", "notes": "呈现最重要的第一个发现，配合图表"},
        {"type": "content", "title": "主要发现（二）", "source": "研究结果", "notes": "呈现第二个核心发现"},
        {"type": "content", "title": "主要发现（三）", "source": "研究结果", "notes": "呈现第三个核心发现或补充发现"},
        {"type": "content", "title": "讨论与意义", "source": "讨论", "notes": "解释发现的理论意义和实践价值"},
        {"type": "content", "title": "研究局限", "source": "讨论/结论", "notes": "诚实说明研究的局限性"},
        {"type": "content", "title": "未来研究方向", "source": "结论", "notes": "提出2-3个有价值的后续研究问题"},
        {"type": "content", "title": "结论", "source": "结论", "notes": "用3-4句话总结核心贡献"},
        {"type": "final", "title": "致谢与问答", "source": "", "notes": "感谢听众，邀请提问"},
    ]

    # Predefined slide templates for defense mode
    _DEFENSE_TEMPLATE = [
        {"type": "title", "title": "封面", "source": "", "notes": "学位论文题目、姓名、导师、专业、日期"},
        {"type": "content", "title": "目录", "source": "", "notes": "展示答辩报告的整体结构"},
        {"type": "content", "title": "研究背景", "source": "第一章", "notes": "详细阐述现实背景和理论背景"},
        {"type": "content", "title": "研究意义", "source": "第一章", "notes": "分理论意义和实践意义两部分"},
        {"type": "content", "title": "文献综述（一）：国内研究", "source": "第二章", "notes": "梳理国内核心研究脉络"},
        {"type": "content", "title": "文献综述（二）：国外研究", "source": "第二章", "notes": "梳理国际研究进展"},
        {"type": "content", "title": "文献述评与研究空白", "source": "第二章", "notes": "总结现有研究不足，引出本研究定位"},
        {"type": "content", "title": "核心概念界定", "source": "第三章", "notes": "对论文中3-5个核心概念进行操作化定义"},
        {"type": "content", "title": "理论框架", "source": "第三章", "notes": "详细阐述理论视角和分析框架"},
        {"type": "content", "title": "研究设计总览", "source": "第四章", "notes": "展示研究设计路线图"},
        {"type": "content", "title": "研究方法", "source": "第四章", "notes": "详细说明定性与定量方法的选择依据"},
        {"type": "content", "title": "数据来源与样本特征", "source": "第四章", "notes": "展示数据收集过程、样本描述性统计"},
        {"type": "content", "title": "研究工具", "source": "第四章", "notes": "量表、问卷、访谈提纲的信效度说明"},
        {"type": "content", "title": "数据分析方法", "source": "第四章", "notes": "说明具体使用的统计或质性分析方法"},
        {"type": "content", "title": "研究发现（一）", "source": "第五章", "notes": "第一个研究问题的发现，配合详细图表"},
        {"type": "content", "title": "研究发现（二）", "source": "第五章", "notes": "第二个研究问题的发现"},
        {"type": "content", "title": "研究发现（三）", "source": "第五章", "notes": "第三个研究问题的发现"},
        {"type": "content", "title": "研究发现（四）", "source": "第五章", "notes": "补充发现或意外发现"},
        {"type": "content", "title": "讨论（一）：与文献对话", "source": "第六章", "notes": "将发现与文献综述中的研究进行对比"},
        {"type": "content", "title": "讨论（二）：理论贡献", "source": "第六章", "notes": "阐述对现有理论的推进或修正"},
        {"type": "content", "title": "讨论（三）：实践启示", "source": "第六章", "notes": "对政策、实践的具体建议"},
        {"type": "content", "title": "研究创新点", "source": "第六章/结论", "notes": "明确列出2-3个创新点"},
        {"type": "content", "title": "研究局限", "source": "第七章", "notes": "从方法、数据、理论三方面说明局限"},
        {"type": "content", "title": "未来研究展望", "source": "第七章", "notes": "提出具体可操作的未来研究方向"},
        {"type": "content", "title": "研究结论", "source": "第七章", "notes": "用 bullet points 总结核心结论"},
        {"type": "final", "title": "致谢", "source": "", "notes": "感谢导师、答辩委员会、家人、同学"},
    ]

    def generate_structure(self, args: dict) -> Dict[str, Any]:
        """Generate slide structure plan.

        Args:
            paper_title: str
            paper_abstract: str
            mode: str - "conference" | "defense"
            key_findings: List[str]
            duration_minutes: int (optional)
        """
        paper_title = args.get("paper_title", "（论文标题）")
        mode = args.get("mode", "conference")
        key_findings = args.get("key_findings", [])
        duration = args.get("duration_minutes", 15 if mode == "conference" else 30)

        template = list(self._CONFERENCE_TEMPLATE if mode == "conference" else self._DEFENSE_TEMPLATE)

        slides = []
        for i, slide_def in enumerate(template, 1):
            slide = {
                "slide_number": i,
                "type": slide_def["type"],
                "title": slide_def["title"],
                "content_bullets": [],
                "speaker_notes": slide_def["notes"],
                "suggested_visual": "",
                "source_section": slide_def["source"],
            }

            # Auto-populate content based on slide type
            if slide["type"] == "title":
                slide["content_bullets"] = [paper_title]
                slide["suggested_visual"] = "无"
            elif "发现" in slide["title"] and key_findings:
                idx = len([s for s in slides if "发现" in s["title"] or "发现" in s.get("title", "")])
                if idx < len(key_findings):
                    slide["content_bullets"] = [key_findings[idx]]
                slide["suggested_visual"] = "图表/数据可视化"
            elif "理论框架" in slide["title"]:
                slide["suggested_visual"] = "概念模型图"
            elif "方法" in slide["title"] or "设计" in slide["title"]:
                slide["suggested_visual"] = "流程图/路线图"
            elif "数据" in slide["title"] or "样本" in slide["title"]:
                slide["suggested_visual"] = "描述性统计表/饼图"
            else:
                slide["suggested_visual"] = "文字要点为主"

            slides.append(slide)

        # Time allocation
        time_per_slide = duration / len(slides)

        return {
            "paper_title": paper_title,
            "mode": mode,
            "duration_minutes": duration,
            "total_slides": len(slides),
            "estimated_time_per_slide": round(time_per_slide, 1),
            "slides": slides,
            "tips": self._get_presentation_tips(mode),
        }

    @staticmethod
    def _get_presentation_tips(mode: str) -> List[str]:
        common = [
            "每页不超过6行文字，每行不超过15字",
            "字号建议：标题32-40pt，正文24-28pt",
            "图表比文字更直观，优先用可视化呈现数据",
            "演讲时面向听众，不要背对听众读PPT",
        ]
        if mode == "conference":
            return common + [
                "会议报告时间紧张，聚焦核心发现",
                "开场30秒内讲清研究问题",
                "结论部分重申贡献",
                "预留2-3分钟用于问答",
            ]
        else:
            return common + [
                "答辩报告更详细，须展示研究过程的严谨性",
                "重点突出创新点和理论贡献",
                "对研究局限要诚实但不过度贬低",
                "预判评委可能提出的问题并准备",
                "控制语速，给评委记录和思考的时间",
            ]


class HTMLSlideRenderer:
    """Render slides as a self-contained HTML file."""

    _HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: "Microsoft YaHei", "PingFang SC", "Hiragino Sans GB", sans-serif; background: #1a1a2e; overflow: hidden; }}
  .slide {{ display: none; width: 100vw; height: 100vh; padding: 60px 80px; flex-direction: column; justify-content: center; align-items: center; position: relative; }}
  .slide.active {{ display: flex; }}
  .slide-title {{ background: linear-gradient(135deg, #16213e 0%, #0f3460 100%); color: #e94560; }}
  .slide-content {{ background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); color: #212529; }}
  .slide-final {{ background: linear-gradient(135deg, #16213e 0%, #1a1a2e 100%); color: #fff; }}
  h1 {{ font-size: 48px; margin-bottom: 24px; text-align: center; line-height: 1.4; }}
  h2 {{ font-size: 40px; margin-bottom: 40px; color: #0f3460; border-bottom: 3px solid #e94560; padding-bottom: 12px; width: 100%; max-width: 900px; }}
  .slide-final h2 {{ color: #e94560; border-bottom-color: #fff; }}
  ul {{ list-style: none; width: 100%; max-width: 900px; }}
  li {{ font-size: 26px; line-height: 1.8; margin-bottom: 16px; padding-left: 36px; position: relative; }}
  li::before {{ content: "▸"; position: absolute; left: 0; color: #e94560; font-weight: bold; }}
  .slide-title h1 {{ color: #fff; font-size: 56px; }}
  .slide-title p {{ color: #a0a0a0; font-size: 24px; margin-top: 20px; }}
  .notes {{ position: absolute; bottom: 30px; left: 80px; right: 80px; font-size: 16px; color: #6c757d; background: rgba(255,255,255,0.9); padding: 12px 20px; border-radius: 8px; border-left: 4px solid #e94560; }}
  .slide-title .notes {{ background: rgba(0,0,0,0.5); color: #ccc; }}
  .slide-final .notes {{ background: rgba(0,0,0,0.4); color: #ccc; }}
  .slide-number {{ position: absolute; bottom: 30px; right: 40px; font-size: 18px; color: #adb5bd; }}
  .nav-hint {{ position: fixed; bottom: 10px; left: 50%; transform: translateX(-50%); font-size: 14px; color: #6c757d; background: rgba(255,255,255,0.8); padding: 6px 16px; border-radius: 20px; z-index: 100; }}
  .progress {{ position: fixed; top: 0; left: 0; height: 4px; background: #e94560; transition: width 0.3s; z-index: 100; }}
  @media print {{
    .slide {{ display: flex !important; page-break-after: always; height: auto; min-height: 100vh; }}
    .nav-hint, .progress {{ display: none; }}
  }}
</style>
</head>
<body>
<div class="progress" id="progress"></div>
{slides}
<div class="nav-hint">← → 方向键翻页 | 按 F11 全屏演示</div>
<script>
  const slides = document.querySelectorAll('.slide');
  let current = 0;
  function show(idx) {{
    slides.forEach(s => s.classList.remove('active'));
    slides[idx].classList.add('active');
    document.getElementById('progress').style.width = ((idx + 1) / slides.length * 100) + '%';
  }}
  document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{ if (current < slides.length - 1) show(++current); }}
    if (e.key === 'ArrowLeft') {{ if (current > 0) show(--current); }}
  }});
  show(0);
</script>
</body>
</html>"""

    def render(self, args: dict) -> Dict[str, Any]:
        """Render slides to HTML.

        Args:
            slides: List[Dict] - from generate_structure
            title: str
        """
        slides_data = args.get("slides", [])
        title = args.get("title", "学术汇报")

        if not slides_data:
            return {"error": "No slides data provided"}

        slide_html_parts = []
        for slide in slides_data:
            slide_type = slide.get("type", "content")
            slide_number = slide.get("slide_number", 0)
            slide_title = html.escape(slide.get("title", ""))
            bullets = slide.get("content_bullets", [])
            notes = html.escape(slide.get("speaker_notes", ""))
            visual = html.escape(slide.get("suggested_visual", ""))
            source = html.escape(slide.get("source_section", ""))

            css_class = {
                "title": "slide-title",
                "final": "slide-final",
            }.get(slide_type, "slide-content")

            bullet_html = ""
            if bullets:
                items = "".join(f"<li>{html.escape(b)}</li>" for b in bullets)
                bullet_html = f"<ul>{items}</ul>"
            else:
                bullet_html = '<ul><li>（请在此补充具体内容）</li></ul>'

            notes_html = ""
            if notes:
                notes_html = f'<div class="notes">💡 演讲备注：{notes}'
                if visual:
                    notes_html += f' | 建议图表：{visual}'
                if source:
                    notes_html += f' | 内容来源：{source}'
                notes_html += '</div>'

            html_part = f'''<div class="slide {css_class}">
  <h2>{slide_title}</h2>
  {bullet_html}
  {notes_html}
  <div class="slide-number">{slide_number}/{len(slides_data)}</div>
</div>'''
            slide_html_parts.append(html_part)

        full_html = self._HTML_TEMPLATE.format(
            title=html.escape(title),
            slides="\n".join(slide_html_parts),
        )

        return {
            "html": full_html,
            "slide_count": len(slides_data),
            "title": title,
            "note": "HTML文件可直接在浏览器中打开，按F11进入全屏演示模式，使用左右方向键翻页。",
        }


class PPTXSlideRenderer:
    """Render slides as a real PPTX file."""
    
    def render(self, args: dict) -> Dict[str, Any]:
        """Render slides to a real PPTX file using python-pptx.
        
        Args:
            slides (List[Dict]): Slides from generate_structure
            title (str): Presentation title
            output_path (str): File path to save
            image_paths (Dict[int, str]): Map of slide index to image path (e.g. {1: "/tmp/plot.png"})
        """
        if not PPTX_AVAILABLE:
            return {"error": "python-pptx package is not installed. Add it to dependencies."}
            
        slides_data = args.get("slides", [])
        title = args.get("title", "Presentation")
        output_path = args.get("output_path", "presentation.pptx")
        image_paths = args.get("image_paths", {})
        
        # Make int keys out of str keys if any
        images = {}
        for k, v in image_paths.items():
            try:
                images[int(k)] = v
            except ValueError:
                pass
                
        prs = Presentation()
        
        # Title slide layout
        title_slide_layout = prs.slide_layouts[0]
        # Bullet slide layout
        bullet_slide_layout = prs.slide_layouts[1]
        
        for idx, slide_data in enumerate(slides_data):
            slide_type = slide_data.get("type", "content")
            slide_title_text = slide_data.get("title", f"Slide {idx+1}")
            bullets = slide_data.get("content_bullets", [])
            notes = slide_data.get("speaker_notes", "")
            
            if slide_type == "title" or slide_type == "final":
                slide = prs.slides.add_slide(title_slide_layout)
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]
                
                title_shape.text = slide_title_text
                
                if bullets:
                    subtitle_shape.text = "
".join(bullets)
                else:
                    subtitle_shape.text = title if slide_type == "title" else "Thank You"
            else:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_title_text
                
                tf = body_shape.text_frame
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
            
            # Add images if provided for this slide
            if idx in images and images[idx]:
                import os
                if os.path.exists(images[idx]):
                    # Add picture (left=1 inch, top=2 inches, width=5 inches)
                    # Adjust parameters dynamically based on slide
                    try:
                        slide.shapes.add_picture(images[idx], Inches(2), Inches(2), width=Inches(6))
                    except Exception as e:
                        print(f"Failed to add image to slide {idx}: {e}")
                        
            # Add speaker notes
            if hasattr(slide, "notes_slide") and notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes
                
        prs.save(output_path)
        
        return {
            "status": "success",
            "format": "pptx",
            "path": output_path,
            "slide_count": len(slides_data),
            "images_embedded": len(images)
        }
